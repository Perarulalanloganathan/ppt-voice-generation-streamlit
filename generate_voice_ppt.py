#!/usr/bin/env python3
"""
PPT Voice Generation Tool

Flow:
  1. Extract text from each slide (python-pptx)
  2. Generate spoken voiceover script per slide (Draup LLM)
  3. Convert each script to MP3 (gTTS)
  4. Embed each MP3 into its slide with auto-play (OOXML / ZIP manipulation)

Usage:
  python generate_voice_ppt.py                          # auto-detects .pptx in CWD
  python generate_voice_ppt.py input.pptx               # specific input
  python generate_voice_ppt.py input.pptx output.pptx   # specify both

Requirements:
    Configure Draup LLM credentials required by your environment.
    Optional overrides:
        DRAUP_LLM_ENV (default: dev)
        DRAUP_LLM_USER (default: nemili)
        DRAUP_LLM_PROVIDER (default: gemini)
        DRAUP_LLM_PROCESS (default: ppt_voice_generation)
        DRAUP_LLM_MODEL (default: gemini/gemini-2.5-flash-lite)
        DRAUP_INPUT_PRICE_PER_MILLION (default: 0.10)
        DRAUP_OUTPUT_PRICE_PER_MILLION (default: 0.40)
        PPT_AUDIO_EMBED_METHOD (default: auto; values: auto|com|ooxml)
    Install: pip install draup_packages python-pptx gTTS lxml python-dotenv imageio-ffmpeg
    Optional on Windows (recommended for stable audio embedding): pip install pywin32
    Note: audio speed adjustment requires FFmpeg. This script auto-uses imageio-ffmpeg if system FFmpeg is missing.
"""

import os
import re
import sys
import shutil
import subprocess
from importlib import import_module
import tempfile
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from gtts import gTTS
from dotenv import load_dotenv

# ── OOXML Namespaces ──────────────────────────────────────────────────────────
P_NS   = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS   = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
REL_NS = 'http://schemas.openxmlformats.org/package/2006/relationships'
CT_NS  = 'http://schemas.openxmlformats.org/package/2006/content-types'
P14_NS = 'http://schemas.microsoft.com/office/powerpoint/2010/main'

AUDIO_REL_TYPE = f'{R_NS}/audio'
MEDIA_REL_TYPE = 'http://schemas.microsoft.com/office/2007/relationships/media'


def _create_draup_llm_handler():
    """Create a Draup LLM handler using notebook-style defaults."""
    try:
        from draup_packages.draup_llm_manager import DraupLLMManager
    except ImportError as exc:
        raise RuntimeError(
            "draup_packages is not installed. Install it before running this script."
        ) from exc

    return DraupLLMManager(
        env=os.environ.get('DRAUP_LLM_ENV', 'dev'),
        user=os.environ.get('DRAUP_LLM_USER', 'nemili'),
        llm_provider=os.environ.get('DRAUP_LLM_PROVIDER', 'gemini'),
        process=os.environ.get('DRAUP_LLM_PROCESS', 'ppt_voice_generation'),
    )


def _extract_draup_content_text(content: Any) -> str:
    """Normalize Draup message content into a plain text string."""
    if isinstance(content, str):
        return content.strip()

    if isinstance(content, list):
        chunks = []
        for item in content:
            text = None
            if isinstance(item, dict):
                text = item.get('text')
            else:
                text = getattr(item, 'text', None)
            if text:
                chunks.append(str(text).strip())
        return '\n'.join([c for c in chunks if c]).strip()

    return str(content).strip()


def _extract_token_usage(resp: Any) -> tuple[int, int, int]:
    """Extract token usage from Draup response in a defensive way."""
    usage = getattr(resp, 'usage', None)
    if usage is None:
        return 0, 0, 0

    input_tokens = int(getattr(usage, 'prompt_tokens', 0) or 0)
    output_tokens = int(getattr(usage, 'completion_tokens', 0) or 0)
    total_tokens = int(getattr(usage, 'total_tokens', 0) or 0)

    if total_tokens == 0:
        total_tokens = input_tokens + output_tokens

    return input_tokens, output_tokens, total_tokens


# ── Step 1: Extract slide texts ───────────────────────────────────────────────
def extract_slide_texts(pptx_path: str) -> list[str]:
    print("Step 1: Extracting slide text...")
    prs = Presentation(pptx_path)
    texts = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
        text = '\n'.join(parts)
        texts.append(text)
        print(f"  Slide {i + 1}: {len(text)} chars")
    print(f"  -> {len(texts)} slides found\n")
    return texts


# ── Step 2: Generate voiceover scripts via Draup LLM ──────────────────────────
def generate_voiceovers(slide_texts: list[str], llm_handler, model: str) -> tuple[list[str], dict[str, float]]:
    print("Step 2: Generating voiceover scripts (Draup LLM)...")
    voiceovers = []
    total_input_tokens = 0
    total_output_tokens = 0
    total_tokens = 0

    input_price_per_million = float(os.environ.get('DRAUP_INPUT_PRICE_PER_MILLION', '0.10'))
    output_price_per_million = float(os.environ.get('DRAUP_OUTPUT_PRICE_PER_MILLION', '0.40'))

    system_prompt = (
        "You are a presentation narrator. Produce concise, natural spoken narration "
        "for slide content. Keep the voice conversational and clear for a live audience."
    )

    for i, text in enumerate(slide_texts):
        if not text.strip():
            voiceovers.append(f"Slide {i + 1}.")
            print(f"  Slide {i + 1}: (empty — placeholder used)")
            continue

        print(f"  Slide {i + 1}...", end=" ", flush=True)
        resp = llm_handler.completion(
            model=model,
            messages=[{
                "role": "system",
                "content": system_prompt,
            }, {
                "role": "user",
                "content": (
                    "You are a presentation narrator. Write a concise, natural spoken "
                    "voiceover (2–4 sentences) for the following slide content. "
                    "Write it conversationally, as if explaining to an audience. "
                    "Do NOT start with 'This slide' or 'In this slide'.\n\n"
                    f"Slide content:\n{text}\n\nVoiceover:"
                )
            }]
        )

        try:
            content = resp.choices[0].message.content
        except (AttributeError, IndexError, TypeError) as exc:
            raise RuntimeError(f"Unexpected Draup response format on slide {i + 1}") from exc

        input_tokens, output_tokens, slide_total_tokens = _extract_token_usage(resp)
        total_input_tokens += input_tokens
        total_output_tokens += output_tokens
        total_tokens += slide_total_tokens

        vo = _extract_draup_content_text(content)
        if not vo:
            vo = f"Slide {i + 1}."

        voiceovers.append(vo)
        preview = vo[:75] + ('...' if len(vo) > 75 else '')
        print(f'"{preview}" (tokens: in={input_tokens}, out={output_tokens}, total={slide_total_tokens})')

    print()
    input_cost = input_price_per_million * (total_input_tokens / 1_000_000)
    output_cost = output_price_per_million * (total_output_tokens / 1_000_000)
    total_cost = input_cost + output_cost

    token_summary = {
        'input_tokens': float(total_input_tokens),
        'output_tokens': float(total_output_tokens),
        'total_tokens': float(total_tokens),
        'input_cost': input_cost,
        'output_cost': output_cost,
        'total_cost': total_cost,
        'input_price_per_million': input_price_per_million,
        'output_price_per_million': output_price_per_million,
    }
    return voiceovers, token_summary


# ── Step 3: Text → MP3 via gTTS ───────────────────────────────────────────────
def _normalize_tts_text(text: str) -> str:
    """
    Normalize narration text to reduce unnatural pauses in TTS playback.
    This helps produce a more natural 1x listening experience.
    """
    cleaned = text.replace('\r\n', '\n').replace('\r', '\n')
    cleaned = re.sub(r'\n+', '. ', cleaned)
    cleaned = re.sub(r'\s+', ' ', cleaned).strip()

    if cleaned and cleaned[-1] not in '.!?':
        cleaned += '.'

    return cleaned


def _get_audio_speed() -> float:
    """Get playback speed for generated narration audio."""
    raw = os.environ.get('PPT_AUDIO_SPEED', '1.5').strip()
    try:
        speed = float(raw)
    except ValueError as exc:
        raise ValueError(f"Invalid PPT_AUDIO_SPEED value: {raw!r}") from exc

    if speed <= 0:
        raise ValueError(f"PPT_AUDIO_SPEED must be > 0, got {speed}")

    return speed


def _apply_audio_speed(path: str, speed: float) -> None:
    """Apply playback speed to an MP3 file in place."""
    if abs(speed - 1.0) < 1e-9:
        return

    ffmpeg_exe = shutil.which('ffmpeg')
    if ffmpeg_exe is None:
        try:
            imageio_ffmpeg = import_module('imageio_ffmpeg')
        except ImportError as exc:
            raise RuntimeError(
                "FFmpeg is required for MP3 speed changes. Install system ffmpeg or: "
                "pip install imageio-ffmpeg"
            ) from exc
        ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()

    # ffmpeg's atempo accepts 0.5..2.0, so split into chained filters if needed.
    factors: list[float] = []
    remaining = speed
    while remaining > 2.0:
        factors.append(2.0)
        remaining /= 2.0
    while remaining < 0.5:
        factors.append(0.5)
        remaining /= 0.5
    factors.append(remaining)
    atempo = ','.join(f'atempo={factor:.6f}' for factor in factors)

    tmp_path = f"{path}.speedtmp.mp3"
    cmd = [
        ffmpeg_exe,
        '-y',
        '-i', path,
        '-filter:a', atempo,
        '-vn',
        tmp_path,
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        stderr_tail = (proc.stderr or '').strip().splitlines()
        details = '\n'.join(stderr_tail[-8:]) if stderr_tail else 'Unknown ffmpeg error.'
        raise RuntimeError(f"Failed to apply audio speed with ffmpeg:\n{details}")

    os.replace(tmp_path, path)


def generate_audio_files(voiceovers: list[str], out_dir: str) -> list[str]:
    audio_speed = _get_audio_speed()
    print(f"Step 3: Converting voiceovers to audio (gTTS, {audio_speed:.2f}x)...")
    paths = []
    for i, text in enumerate(voiceovers):
        path = os.path.join(out_dir, f"slide_{i + 1:02d}.mp3")
        print(f"  Slide {i + 1}...")
        tts_text = _normalize_tts_text(text)
        gTTS(text=tts_text, lang='en', slow=False).save(path)
        _apply_audio_speed(path, audio_speed)
        paths.append(path)
    print()
    return paths


# ── Step 4: Embed audio into PPTX ─────────────────────────────────────────────
def _get_next_rel_id(rels_root) -> str:
    """Find the next unused rId number."""
    used = set()
    for rel in rels_root.findall(f'{{{REL_NS}}}Relationship'):
        rid = rel.get('Id', '')
        if rid.startswith('rId'):
            try:
                used.add(int(rid[3:]))
            except ValueError:
                pass
    return f'rId{max(used, default=0) + 1}'


def _get_next_shape_id(slide_root) -> int:
    """Find the highest existing shape id and return id+1."""
    used = set()
    for el in slide_root.iter():
        v = el.get('id')
        if v:
            try:
                used.add(int(v))
            except ValueError:
                pass
    return max(used, default=0) + 1


def _make_audio_shape(shape_id: int, audio_link_r_id: str, media_embed_r_id: str):
    """
    Build a hidden audio shape (positioned off-screen at -1in, -1in).
    The shape references the audio file via <a:audioFile r:link="rId">.
    """
    sp = etree.Element(f'{{{P_NS}}}sp')

    nvSpPr = etree.SubElement(sp, f'{{{P_NS}}}nvSpPr')
    etree.SubElement(nvSpPr, f'{{{P_NS}}}cNvPr',
                     id=str(shape_id), name=f'AudioShape_{shape_id}')
    etree.SubElement(nvSpPr, f'{{{P_NS}}}cNvSpPr')
    nvPr = etree.SubElement(nvSpPr, f'{{{P_NS}}}nvPr')
    af = etree.SubElement(nvPr, f'{{{A_NS}}}audioFile')
    af.set(f'{{{R_NS}}}link', audio_link_r_id)

    # PowerPoint expects an Office 2010 media extension that embeds the media rel.
    ext_lst = etree.SubElement(nvPr, f'{{{P_NS}}}extLst')
    ext = etree.SubElement(ext_lst, f'{{{P_NS}}}ext')
    ext.set('uri', '{DAA4B4D4-2A88-4E4E-A5AB-9C7D1A2902A8}')
    media = etree.SubElement(ext, f'{{{P14_NS}}}media')
    media.set(f'{{{R_NS}}}embed', media_embed_r_id)

    spPr = etree.SubElement(sp, f'{{{P_NS}}}spPr')
    xfrm = etree.SubElement(spPr, f'{{{A_NS}}}xfrm')
    etree.SubElement(xfrm, f'{{{A_NS}}}off', x='-914400', y='-914400')   # off-screen
    etree.SubElement(xfrm, f'{{{A_NS}}}ext', cx='914400', cy='914400')
    prstGeom = etree.SubElement(spPr, f'{{{A_NS}}}prstGeom', prst='rect')
    etree.SubElement(prstGeom, f'{{{A_NS}}}avLst')

    return sp


def _make_timing(shape_id: int):
    """
    Build a <p:timing> element that auto-plays the audio when the slide opens.
    delay="0" on the outer stCondLst means "start immediately" (not on click).
    """
    timing = etree.Element(f'{{{P_NS}}}timing')
    tnLst  = etree.SubElement(timing, f'{{{P_NS}}}tnLst')
    par0   = etree.SubElement(tnLst, f'{{{P_NS}}}par')

    cTn1 = etree.SubElement(par0, f'{{{P_NS}}}cTn',
                             id='1', dur='indefinite', restart='never', nodeType='tmRoot')
    cl1  = etree.SubElement(cTn1, f'{{{P_NS}}}childTnLst')

    seq  = etree.SubElement(cl1, f'{{{P_NS}}}seq', concurrent='1', nextAc='seek')
    cTn2 = etree.SubElement(seq, f'{{{P_NS}}}cTn',
                             id='2', dur='indefinite', nodeType='mainSeq')
    cl2  = etree.SubElement(cTn2, f'{{{P_NS}}}childTnLst')

    # Outer animation group — delay="0" = auto-play on slide open
    par3 = etree.SubElement(cl2, f'{{{P_NS}}}par')
    cTn3 = etree.SubElement(par3, f'{{{P_NS}}}cTn', id='3', fill='hold')
    sc3  = etree.SubElement(cTn3, f'{{{P_NS}}}stCondLst')
    etree.SubElement(sc3, f'{{{P_NS}}}cond', delay='0')
    cl3  = etree.SubElement(cTn3, f'{{{P_NS}}}childTnLst')

    # Inner animation group
    par4 = etree.SubElement(cl3, f'{{{P_NS}}}par')
    cTn4 = etree.SubElement(par4, f'{{{P_NS}}}cTn', id='4', fill='hold')
    sc4  = etree.SubElement(cTn4, f'{{{P_NS}}}stCondLst')
    etree.SubElement(sc4, f'{{{P_NS}}}cond', delay='0')
    cl4  = etree.SubElement(cTn4, f'{{{P_NS}}}childTnLst')

    # Audio playback node targeting our shape
    audio_el = etree.SubElement(cl4, f'{{{P_NS}}}audio')
    cMedia   = etree.SubElement(audio_el, f'{{{P_NS}}}cMediaNode',
                                vol='80000', mute='0', numSld='0', showWhenStopped='1')
    etree.SubElement(cMedia, f'{{{P_NS}}}cTn', id='5', fill='hold')
    tgtEl = etree.SubElement(cMedia, f'{{{P_NS}}}tgtEl')
    etree.SubElement(tgtEl, f'{{{P_NS}}}spTgt', spid=str(shape_id))

    # prev / next navigation conditions
    prevCL = etree.SubElement(seq, f'{{{P_NS}}}prevCondLst')
    prevC  = etree.SubElement(prevCL, f'{{{P_NS}}}cond', evt='onPrev', delay='0')
    etree.SubElement(prevC, f'{{{P_NS}}}tn')

    nextCL = etree.SubElement(seq, f'{{{P_NS}}}nextCondLst')
    nextC  = etree.SubElement(nextCL, f'{{{P_NS}}}cond', evt='onNext', delay='0')
    etree.SubElement(nextC, f'{{{P_NS}}}tn')

    etree.SubElement(timing, f'{{{P_NS}}}bldLst')
    return timing


def _ensure_mp3_content_type(tmp_dir: str):
    """Add MP3 content type to [Content_Types].xml if not already present."""
    ct_path = os.path.join(tmp_dir, '[Content_Types].xml')
    tree = etree.parse(ct_path)
    root = tree.getroot()
    for d in root.findall(f'{{{CT_NS}}}Default'):
        if d.get('Extension', '').lower() == 'mp3':
            return
    d = etree.SubElement(root, f'{{{CT_NS}}}Default')
    d.set('Extension', 'mp3')
    d.set('ContentType', 'audio/mpeg')
    with open(ct_path, 'wb') as f:
        f.write(etree.tostring(root, xml_declaration=True,
                               encoding='UTF-8', standalone=True))


def embed_audio_into_pptx(pptx_path: str, audio_paths: list[str], output_path: str):
    print("Step 4: Embedding audio into PPTX...")
    tmp = tempfile.mkdtemp()
    try:
        # Unpack the PPTX (it's a ZIP)
        with zipfile.ZipFile(pptx_path, 'r') as z:
            z.extractall(tmp)

        slides_dir = os.path.join(tmp, 'ppt', 'slides')
        rels_dir   = os.path.join(slides_dir, '_rels')
        media_dir  = os.path.join(tmp, 'ppt', 'media')
        os.makedirs(media_dir, exist_ok=True)
        os.makedirs(rels_dir, exist_ok=True)

        # Sorted slide XML files (slide1.xml, slide2.xml, ...)
        slide_files = sorted(
            [f for f in os.listdir(slides_dir)
             if f.startswith('slide') and f.endswith('.xml')],
            key=lambda x: int(''.join(filter(str.isdigit, x)) or '0')
        )

        for i, (slide_file, audio_path) in enumerate(zip(slide_files, audio_paths)):
            if not os.path.isfile(audio_path):
                print(f"  Slide {i + 1}: audio file missing — skipped")
                continue

            print(f"  Slide {i + 1} ({slide_file})")

            # 1. Copy MP3 into ppt/media/
            audio_fname = f'audio_{i + 1:02d}.mp3'
            shutil.copy2(audio_path, os.path.join(media_dir, audio_fname))

            # 2. Add audio/media relationships to slide .rels file
            rels_path = os.path.join(rels_dir, f'{slide_file}.rels')
            if os.path.exists(rels_path):
                rels_root = etree.parse(rels_path).getroot()
            else:
                rels_root = etree.Element(f'{{{REL_NS}}}Relationships')

            audio_r_id = _get_next_rel_id(rels_root)
            audio_rel = etree.SubElement(rels_root, f'{{{REL_NS}}}Relationship')
            audio_rel.set('Id', audio_r_id)
            audio_rel.set('Type', AUDIO_REL_TYPE)
            audio_rel.set('Target', f'../media/{audio_fname}')

            media_r_id = _get_next_rel_id(rels_root)
            media_rel = etree.SubElement(rels_root, f'{{{REL_NS}}}Relationship')
            media_rel.set('Id', media_r_id)
            media_rel.set('Type', MEDIA_REL_TYPE)
            media_rel.set('Target', f'../media/{audio_fname}')

            with open(rels_path, 'wb') as f:
                f.write(etree.tostring(rels_root, xml_declaration=True,
                                       encoding='UTF-8', standalone=True))

            # 3. Add hidden audio shape to the slide's spTree
            slide_path = os.path.join(slides_dir, slide_file)
            slide_root = etree.parse(slide_path).getroot()

            sp_tree = slide_root.find(f'.//{{{P_NS}}}spTree')
            if sp_tree is None:
                print(f"    WARNING: no spTree found — skipped")
                continue

            shape_id = _get_next_shape_id(slide_root)
            sp_tree.append(_make_audio_shape(shape_id, audio_r_id, media_r_id))

            # 4. Replace slide timing with auto-play timing
            for old_timing in slide_root.findall(f'{{{P_NS}}}timing'):
                slide_root.remove(old_timing)
            slide_root.append(_make_timing(shape_id))

            with open(slide_path, 'wb') as f:
                f.write(etree.tostring(slide_root, xml_declaration=True,
                                       encoding='UTF-8', standalone=True))

        # 5. Register MP3 content type
        _ensure_mp3_content_type(tmp)

        # 6. Re-pack as PPTX
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as z:
            for dirpath, _, filenames in os.walk(tmp):
                for fname in filenames:
                    fpath   = os.path.join(dirpath, fname)
                    arcname = os.path.relpath(fpath, tmp)
                    z.write(fpath, arcname)

        print(f"\n  Saved -> {output_path}")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def embed_audio_with_powerpoint(pptx_path: str, audio_paths: list[str], output_path: str):
    """Embed audio using native PowerPoint automation on Windows."""
    print("Step 4: Embedding audio into PPTX (PowerPoint COM)...")
    try:
        import pythoncom
        import win32com.client
    except ImportError as exc:
        raise RuntimeError(
            "pywin32 is required for PowerPoint COM embedding on Windows. "
            "Install it with: pip install pywin32"
        ) from exc

    abs_input = str(Path(pptx_path).resolve())
    abs_output = str(Path(output_path).resolve())
    abs_audio_paths = [str(Path(p).resolve()) for p in audio_paths]

    pythoncom.CoInitialize()
    app = None
    presentation = None
    try:
        app = win32com.client.DispatchEx('PowerPoint.Application')
        presentation = app.Presentations.Open(abs_input, WithWindow=False)

        slide_count = presentation.Slides.Count
        for i, audio_path in enumerate(abs_audio_paths):
            if i + 1 > slide_count:
                break
            if not os.path.isfile(audio_path):
                print(f"  Slide {i + 1}: audio file missing — skipped")
                continue

            print(f"  Slide {i + 1}...")
            slide = presentation.Slides(i + 1)

            try:
                media_shape = slide.Shapes.AddMediaObject2(
                    FileName=audio_path,
                    LinkToFile=False,
                    SaveWithDocument=True,
                    Left=-100.0,
                    Top=-100.0,
                    Width=1.0,
                    Height=1.0,
                )
            except Exception:
                media_shape = slide.Shapes.AddMediaObject(
                    FileName=audio_path,
                    LinkToFile=False,
                    SaveWithDocument=True,
                    Left=-100.0,
                    Top=-100.0,
                    Width=1.0,
                    Height=1.0,
                )

            play_settings = media_shape.AnimationSettings.PlaySettings
            play_settings.PlayOnEntry = True
            play_settings.HideWhileNotPlaying = True
            play_settings.StopAfterSlides = 1

        presentation.SaveCopyAs(abs_output)
        print(f"\n  Saved -> {output_path}")
    finally:
        if presentation is not None:
            presentation.Close()
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def embed_audio(pptx_path: str, audio_paths: list[str], output_path: str):
    """
    Embed audio with a robust strategy:
    - On Windows, prefer native PowerPoint COM embedding for maximum compatibility.
    - Otherwise, use OOXML fallback.
    """
    method = os.environ.get('PPT_AUDIO_EMBED_METHOD', 'auto').strip().lower()

    if method in ('auto', 'com') and sys.platform.startswith('win'):
        try:
            embed_audio_with_powerpoint(pptx_path, audio_paths, output_path)
            return
        except Exception as exc:
            if method == 'com':
                raise
            print(f"WARNING: PowerPoint COM embedding failed ({exc}). Falling back to OOXML.")

    embed_audio_into_pptx(pptx_path, audio_paths, output_path)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    # Load environment variables from .env so credentials are available.
    load_dotenv()

    # Resolve input path
    if len(sys.argv) >= 2:
        input_path = sys.argv[1]
    else:
        pptx_files = sorted(Path('.').glob('*.pptx'))
        if not pptx_files:
            print("No .pptx file found in current directory.")
            print("Usage: python generate_voice_ppt.py input.pptx [output.pptx]")
            sys.exit(1)
        input_path = str(pptx_files[0])
        print(f"Auto-detected: {input_path}")

    # Resolve output path
    if len(sys.argv) >= 3:
        output_path = sys.argv[2]
    else:
        stem = Path(input_path).stem
        output_path = str(Path(input_path).parent / f"{stem}_with_audio.pptx")

    print(f"\n{'=' * 60}")
    print(f"  PPT Voice Generation")
    print(f"  Input : {Path(input_path).name}")
    print(f"  Output: {Path(output_path).name}")
    print(f"{'=' * 60}\n")

    # Initialize Draup LLM handler
    try:
        llm_handler = _create_draup_llm_handler()
    except Exception as exc:
        print(f"ERROR: failed to initialize Draup LLM: {exc}")
        sys.exit(1)

    llm_model = os.environ.get('DRAUP_LLM_MODEL', 'gemini/gemini-2.5-flash-lite')
    print(f"Using Draup model: {llm_model}")

    tmp_audio = tempfile.mkdtemp()
    try:
        texts       = extract_slide_texts(input_path)
        voiceovers, token_summary = generate_voiceovers(texts, llm_handler, llm_model)

        print(f"{'=' * 60}")
        print("  LLM Token & Cost Summary")
        print(f"  Input Tokens : {int(token_summary['input_tokens']):,}")
        print(f"  Output Tokens: {int(token_summary['output_tokens']):,}")
        print(f"  Total Tokens : {int(token_summary['total_tokens']):,}")
        print(
            f"  Input Cost   : ${token_summary['input_cost']:.6f} "
            f"(@ ${token_summary['input_price_per_million']:.2f}/1M)"
        )
        print(
            f"  Output Cost  : ${token_summary['output_cost']:.6f} "
            f"(@ ${token_summary['output_price_per_million']:.2f}/1M)"
        )
        print(f"  Total Cost   : ${token_summary['total_cost']:.6f}")
        print(f"{'=' * 60}\n")

        audio_paths = generate_audio_files(voiceovers, tmp_audio)
        embed_audio(input_path, audio_paths, output_path)

        print(f"\n{'=' * 60}")
        print(f"  Done!  ->  {Path(output_path).name}")
        print(f"{'=' * 60}")
    finally:
        shutil.rmtree(tmp_audio, ignore_errors=True)


if __name__ == '__main__':
    main()
